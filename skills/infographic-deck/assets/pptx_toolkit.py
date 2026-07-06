#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Toolkit для инфографических слайдов (редактируемый pptx, шрифты встроены).
Скопировать в рабочий скрипт, задать палитру, вызывать helpers, в конце embed_fonts().
Координаты в пикселях холста 1280x720 (1px = 1/96 inch). Проверено в PowerPoint и LibreOffice.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import zipfile, os

# ---- ПАЛИТРА (задать под проект из palettes.md; НЕ жёлто-бирюзовую) ----
ACCENT="2456E6"; ACCENT_D="1B3FB0"; INK="14203A"; PANEL="EEF2FB"
MUTED="5B6675"; LINE="DCE3EF"; CARD="FFFFFF"; GREY="C4BFB4"
CG="Cabinet Grotesk"; GS="General Sans"
L=PP_ALIGN.LEFT; R=PP_ALIGN.RIGHT; C=PP_ALIGN.CENTER
TOP=MSO_ANCHOR.TOP; MID=MSO_ANCHOR.MIDDLE; BOT=MSO_ANCHOR.BOTTOM

prs=Presentation(); prs.slide_width=Inches(13.3333); prs.slide_height=Inches(7.5)
def new_slide(bg=CARD):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor.from_string(bg); return s
def px(v): return Inches(v/96.0)
def pt(v): return Pt(v*0.75)
def _nsh(sp):
    p=sp._element.spPr
    if p.find(qn('a:effectLst')) is None: p.append(p.makeelement(qn('a:effectLst'),{}))

def text(s,x,y,w,h,paras,align=L,anchor=TOP,space=1.0):
    tb=s.shapes.add_textbox(px(x),px(y),px(w),px(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=anchor
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if space: p.line_spacing=space
        if isinstance(para,dict): para=[para]
        for r in para:
            run=p.add_run(); run.text=r["t"]; f=run.font
            f.size=pt(r["s"]); f.bold=r.get("b",False); f.italic=r.get("i",False)
            f.name=r.get("f",GS); f.color.rgb=RGBColor.from_string(r["c"])
    return tb
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,rad=None):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if rad is not None else MSO_SHAPE.RECTANGLE
    sp=s.shapes.add_shape(shp,px(x),px(y),px(w),px(h)); sp.shadow.inherit=False
    if rad is not None:
        try: sp.adjustments[0]=rad
        except: pass
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb=RGBColor.from_string(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    _nsh(sp); return sp
def oval(s,x,y,d,fill=None,line=None,lw=1.2,label=None,tc=CARD,fs=15,ff=CG):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,px(x),px(y),px(d),px(d)); sp.shadow.inherit=False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb=RGBColor.from_string(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb=RGBColor.from_string(line); sp.line.width=Pt(lw)
    else: sp.line.fill.background()
    _nsh(sp)
    if label is not None: text(s,x,y-1,d,d,[{"t":label,"s":fs,"c":tc,"b":True,"f":ff}],align=C,anchor=MID)
    return sp

# ---- харви-шарик: полный / половина / пустой (bg = цвет фона ячейки) ----
def ball(s,cx,cy,kind,bg=CARD,d=15):
    x=cx-d/2; y=cy-d/2
    if kind=="f": oval(s,x,y,d,fill=ACCENT,line=ACCENT,lw=1.2)
    elif kind=="o": oval(s,x,y,d,fill=CARD,line=GREY,lw=1.4)
    else:
        oval(s,x,y,d,fill=ACCENT)
        rect(s,x+d/2,y-0.5,d/2+0.6,d+1,fill=bg)
        oval(s,x,y,d,fill=None,line=GREY,lw=1.4)

# ---- верхние вкладки-навигация ----
def navtabs(s,tabs,active_idx,y=24,x0=34,total=1212,h=25):
    n=len(tabs); gap=8; tw=(total-gap*(n-1))/n; x=x0
    for i,lb in enumerate(tabs):
        on=(i==active_idx)
        rect(s,x,y,tw,h,fill=(ACCENT if on else "F1EEE8"),rad=0.5)
        text(s,x,y,tw,h,[{"t":lb,"s":11,"c":(CARD if on else MUTED),"b":True,"f":GS}],align=C,anchor=MID)
        x+=tw+gap

def head(s,kicker,title,tsize=27,x=34,y=56,w=1000):
    text(s,x,y-22,w,20,[{"t":kicker,"s":13,"c":ACCENT,"b":True,"f":GS}])
    text(s,x,y,w,110,[{"t":title,"s":tsize,"c":INK,"b":False,"f":CG}],space=1.06)

def big_number(s,x,y,num,label,w=300,size=64):
    text(s,x,y,w,70,[{"t":num,"s":size,"c":INK,"b":False,"f":CG}],anchor=MID)
    text(s,x,y+size+6,w,50,[{"t":label,"s":14,"c":MUTED,"f":GS}],space=1.15)

def highlight_box(s,x,y,w,h,textstr,fs=13):
    rect(s,x,y,w,h,fill=CARD,line=ACCENT,lw=2.2,rad=0.06)
    text(s,x+14,y,w-28,h,[{"t":textstr,"s":fs,"c":INK,"b":True,"f":GS}],align=C,anchor=MID,space=1.15)

def feature(s,x,y,num,title,body,hl=None):
    oval(s,x,y,44,fill=ACCENT,label=str(num),tc=CARD,fs=20)
    text(s,x+62,y-4,320,40,[{"t":title,"s":26,"c":INK,"b":False,"f":CG}])
    text(s,x+62,y+40,320,90,[{"t":body,"s":15,"c":MUTED,"f":GS}],space=1.2)
    if hl: highlight_box(s,x,y+140,320,60,hl)

def footnote(s,src,x=34,y=700,w=900):
    text(s,x,y,w,16,[{"t":src,"s":9.5,"c":MUTED,"i":True,"f":GS}])

# ---- нативный donut (фон областей выключен, иначе LibreOffice рисует серый бокс) ----
def donut(s,x,y,d,cats,vals,colors):
    cd=CategoryChartData(); cd.categories=cats; cd.add_series("s",vals)
    gf=s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,px(x),px(y),px(d),px(d),cd)
    ch=gf.chart; ch.has_legend=False; ch.has_title=False
    ser=ch.plots[0].series[0]
    for i,col in enumerate(colors):
        p=ser.points[i]; p.format.fill.solid(); p.format.fill.fore_color.rgb=RGBColor.from_string(col)
    cs=ch._chartSpace
    def nf(el):
        sp=el.makeelement(qn('c:spPr'),{}); etree.SubElement(sp,qn('a:noFill'))
        ln=etree.SubElement(sp,qn('a:ln')); etree.SubElement(ln,qn('a:noFill')); return sp
    cs.find(qn('c:chart')).addnext(nf(cs))
    pa=cs.find('.//'+qn('c:plotArea'))
    if pa is not None: pa.append(nf(pa))
    return gf

# ---- горизонтальный bar-чарт из прямоугольников ----
def hbars(s,x,y,rows,bw=600,rh=34,maxv=None):
    maxv=maxv or max(v for _,v in rows)
    yy=y
    for name,v in rows:
        text(s,x,yy,220,rh,[{"t":name,"s":13,"c":INK,"f":GS}],anchor=MID)
        rect(s,x+230,yy+rh/2-10,bw*v/maxv,20,fill=ACCENT)
        text(s,x+230+bw+8,yy,60,rh,[{"t":str(v),"s":13,"c":ACCENT,"b":True,"f":CG}],anchor=MID)
        yy+=rh

# ---- встроить шрифты (ttf пути задать) ----
def embed_fonts(out_path, fonts):
    """fonts: list of (typeface, slot, ttf_path); slot in regular/bold/italic"""
    nP="http://schemas.openxmlformats.org/presentationml/2006/main"
    nR="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    nCT="http://schemas.openxmlformats.org/package/2006/content-types"
    nPR="http://schemas.openxmlformats.org/package/2006/relationships"
    z=zipfile.ZipFile(out_path,'r'); D={n:z.read(n) for n in z.namelist()}; z.close()
    ct=etree.fromstring(D['[Content_Types].xml'])
    if not any(e.get('Extension')=='fntdata' for e in ct if e.tag==f'{{{nCT}}}Default'):
        d=etree.SubElement(ct,f'{{{nCT}}}Default'); d.set('Extension','fntdata'); d.set('ContentType','application/x-fontdata')
    rx=etree.fromstring(D['ppt/_rels/presentation.xml.rels']); mx=0
    for e in rx:
        r=e.get('Id','')
        if r.startswith('rId'):
            try: mx=max(mx,int(r[3:]))
            except: pass
    fam={}; rids=[]
    for i,(tf,slot,p) in enumerate(fonts):
        rid='rId%d'%(mx+1+i); rel=etree.SubElement(rx,f'{{{nPR}}}Relationship')
        rel.set('Id',rid); rel.set('Type',nR+'/font'); rel.set('Target','fonts/font%d.fntdata'%(i+1))
        fam.setdefault(tf,[]).append((slot,rid)); rids.append((i,p))
    pr=etree.fromstring(D['ppt/presentation.xml']); pr.set('embedTrueTypeFonts','1')
    efl=etree.SubElement(pr,f'{{{nP}}}embeddedFontLst')
    for tf,slots in fam.items():
        ef=etree.SubElement(efl,f'{{{nP}}}embeddedFont'); fo=etree.SubElement(ef,f'{{{nP}}}font'); fo.set('typeface',tf)
        for slot,rid in slots:
            e=etree.SubElement(ef,f'{{{nP}}}{slot}'); e.set(f'{{{nR}}}id',rid)
    ns=pr.find(f'{{{nP}}}notesSz')
    if ns is not None: pr.remove(efl); ns.addnext(efl)
    D['[Content_Types].xml']=etree.tostring(ct,xml_declaration=True,encoding='UTF-8',standalone=True)
    D['ppt/_rels/presentation.xml.rels']=etree.tostring(rx,xml_declaration=True,encoding='UTF-8',standalone=True)
    D['ppt/presentation.xml']=etree.tostring(pr,xml_declaration=True,encoding='UTF-8',standalone=True)
    for i,p in rids: D['ppt/fonts/font%d.fntdata'%(i+1)]=open(p,'rb').read()
    tmp=out_path+'.tmp'; zo=zipfile.ZipFile(tmp,'w',zipfile.ZIP_DEFLATED)
    for n,b in D.items(): zo.writestr(n,b)
    zo.close(); os.replace(tmp,out_path)

# ---- пример: один слайд со всеми компонентами ----
if __name__=="__main__":
    s=new_slide()
    navtabs(s,["РЕЗЮМЕ","РЫНОК","АНАЛИЗ ЦА","КОНКУРЕНТЫ","ПРОДУКТ","КОМАНДА"],3)
    head(s,"Конкуренты","Свободный угол, где нужны независимость и глубина одновременно")
    rect(s,34,150,360,300,fill=PANEL,rad=0.03)
    big_number(s,60,180,"930M","TAM, весь рынок")
    highlight_box(s,60,300,300,60,"Пересечение не закрывает никто")
    feature(s,430,160,1,"Простота","Заполняешь один раз, остальное делает агент","Экономит около 3 часов")
    donut(s,900,180,150,["A","B","C"],(5,4,3),[ACCENT,ACCENT_D,GREY])
    footnote(s,"Источник: анализ команды.")
    OUT="demo.pptx"; prs.save(OUT)
    FB="."  # путь к папкам шрифтов
    embed_fonts(OUT,[
        (CG,"regular",FB+"/CabinetGrotesk-Extrabold.ttf"),
        (CG,"bold",   FB+"/CabinetGrotesk-Bold.ttf"),
        (GS,"regular",FB+"/GeneralSans-Regular.ttf"),
        (GS,"bold",   FB+"/GeneralSans-Semibold.ttf"),
        (GS,"italic", FB+"/GeneralSans-Italic.ttf"),
    ])
    print("saved",OUT)
